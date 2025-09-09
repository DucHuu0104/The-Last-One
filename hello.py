import streamlit as st
from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_google_genai import GoogleGenerativeAIEmbeddings, ChatGoogleGenerativeAI
from langchain_community.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from docx import Document
from dotenv import load_dotenv
from pptx import Presentation
import tempfile
import os
import google.generativeai as genai
from io import BytesIO
import matplotlib.pyplot as plt
from collections import Counter
import seaborn as sns
from wordcloud import  STOPWORDS

from PIL import Image
import pytesseract
# Load biến môi trường
load_dotenv()
api_key = os.getenv('GEMINI_API_KEY')

if not api_key:
    st.error('Không tìm thấy API key. Hãy kiểm tra file .env')
    st.stop()

genai.configure(api_key=api_key)

# HÀM XỬ LÝ TÀI LIỆU PDF,DOCX, PPTX, TXT
def get_text(docs):
    text = ''
    try:
        for file in docs:
            if file.name.endswith('.pdf'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
                    tmp_file.write(file.read())
                    tmp_file_path = tmp_file.name
                pdf_reader = PyPDFLoader(tmp_file_path)
                for page in pdf_reader.load_and_split():
                    text += page.page_content
                os.unlink(tmp_file_path)

            elif file.name.endswith('.docx'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
                    tmp_file.write(file.read())
                    tmp_file_path = tmp_file.name
                doc = Document(tmp_file_path)
                for para in doc.paragraphs:
                    text += para.text + '\n'
                os.unlink(tmp_file_path)

            elif file.name.endswith('.pptx'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
                    tmp_file.write(file.read())
                    tmp_file_path = tmp_file.name
                prs = Presentation(tmp_file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + '\n'
                os.unlink(tmp_file_path)

            elif file.name.endswith('.txt'):
                content = file.read().decode('utf-8') 
                text += content + '\n'
                
            elif file.name.endswith(('.png', '.jpg', '.jpeg', '.bmp')):
                with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(file.name)[1]) as tmp_file:
                    tmp_file.write(file.read())
                    tmp_file_path = tmp_file.name
                try:
                    img = Image.open(tmp_file_path)
                    text += pytesseract.image_to_string(img, lang='vie+eng') + '\n'
                except Exception as e:
                    st.error(f'Lỗi OCR cho ảnh {file.name}: {str(e)}')
                finally:
                    os.unlink(tmp_file_path)

            else:
                st.error(f'Định dạng không hỗ trợ: {file.name}')
    except Exception as e:
        st.error(f'Lỗi đọc tài liệu: {str(e)}')
    return text
# CHIA NHỎ VĂN BẢN
def get_text_chunk(text):
    try: 
        text_splitter = RecursiveCharacterTextSplitter(chunk_size=10000, chunk_overlap=1000)
        chunks = text_splitter.split_text(text)
        return chunks
    except Exception as e:
        st.error(f'Lỗi chia chunk: {str(e)}')
        return []

# TẠO VÀ LƯU VECTOR DATABASE
def get_vector_store(text_chunks):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model='models/embedding-001')
        vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
        vector_store.save_local('faiss_index')
        st.success('Tài liệu đã được phân tích và lưu xong.')
    except Exception as e:
        st.error(f'Lỗi lưu FAISS index: {str(e)}')

# TẠO CHUỖI QA VỚI PROMPT
def get_conversation_chain():
    prompt_template = """
    Trả lời câu hỏi một cách chi tiết nhất có thể dựa trên ngữ cảnh được cung cấp. Nếu câu trả lời không có trong ngữ cảnh được cung cấp, hãy nói, "Câu trả lời không có trong ngữ cảnh."
    Không cung cấp thông tin sai lệch.

    Ngữ cảnh: {context}
    Câu hỏi: {question}

    Answer:
    """
    try:
        model = ChatGoogleGenerativeAI(model='gemini-2.0-flash', temperature=0.3)
        prompt = PromptTemplate(template=prompt_template, input_variables=['context', 'question'])
        chain = load_qa_chain(model, chain_type='stuff', prompt=prompt)
        return chain
    except Exception as e:
        st.error(f"Lỗi khởi tạo chuỗi QA: {str(e)}")
        return None


# TRẢ LỜI CÂU HỎI NGƯỜI DÙNG VÀ LƯU LỊCH SỬ
def user_input(user_question):
    try:
        embeddings = GoogleGenerativeAIEmbeddings(model='models/embedding-001')
        if not os.path.exists('faiss_index'):
            st.error('Không tìm thấy FAISS index. Hãy tải và phân tích tài liệu trước.')
            return None
        new_db = FAISS.load_local('faiss_index', embeddings, allow_dangerous_deserialization=True)
        docs = new_db.similarity_search(user_question)
        chain = get_conversation_chain()
        if not chain:
            return None
        response = chain(
            {"input_documents": docs, "question": user_question},
            return_only_outputs=True
        )
        answer = response["output_text"]

        # Lưu lịch sử
        if "chat_history" not in st.session_state:
            st.session_state.chat_history = []
        st.session_state.chat_history.append({"question": user_question, "answer": answer})
        return answer
    except Exception as e:
        st.error(f'Lỗi xử lý câu hỏi: {str(e)}')
        return None


def show_statistics():
    st.markdown("## 📊 Thống kê câu hỏi")

    if "chat_history" not in st.session_state or not st.session_state.chat_history:
        st.info("Chưa có dữ liệu để thống kê.")
        return

    questions = [chat["question"] for chat in st.session_state.chat_history]
    st.write(f"Tổng số câu hỏi đã hỏi: **{len(questions)}**")

    #Thống kê từ khóa phổ biến
    st.markdown("### 🧠 Top 10 từ khóa được hỏi nhiều nhất")

    all_words = " ".join(questions).lower().split()
    #không cho các từ lên top 
    stopwords = set(STOPWORDS) | {"câu", "hỏi", "nào", "gì", "là", "cho", "về", "có", "trong", "và", "là", "nhưng", "thì", "với", "cho", "từ", "đến",
    "một", "những", "các", "của", "này", "kia", "đó", "ấy",
    "tôi", "bạn", "chúng", "mình", "em", "anh", "chị", "họ",
    "ông", "bà", "khi", "nếu", "vì", "để", "cũng", "đã", "đang",
    "sẽ", "phải", "còn", "hay", "như", "nên"}

    words_filtered = [word for word in all_words if word not in stopwords and len(word) > 3]
    common_words = Counter(words_filtered)
    top_words = dict(common_words.most_common(10))

    # Biểu đồ bar
    plt.figure(figsize=(6, 4))
    sns.barplot(x=list(top_words.values()), y=list(top_words.keys()), palette='viridis')
    plt.title("Top 10 từ khóa được dùng nhiều nhất")
    plt.xlabel("Tần suất")
    plt.ylabel("Từ khóa")
    st.pyplot(plt)

# XUẤT LỊCH SỬ HỎI ĐÁP RA FILE WORD
def export_chat_history_to_word(chat_history):
    doc = Document()
    doc.add_heading('Lịch sử hỏi đáp ChatBot', level=1)

    for i, chat in enumerate(chat_history, 1):
        doc.add_paragraph(f"Câu hỏi {i}:", style='List Number')
        doc.add_paragraph(chat['question'])
        doc.add_paragraph(f"Câu trả lời {i}:", style='List Number')
        doc.add_paragraph(chat['answer'])
        doc.add_paragraph()

    file_stream = BytesIO()
    doc.save(file_stream)
    file_stream.seek(0)
    return file_stream

# GIAO DIỆN STREAMLIT
def main():
    st.set_page_config(page_title='ChatBot phân tích tài liệu ')
    st.title('ChatBot phân tích tài liệu')

    user_question = st.text_input('📌 Nhập câu hỏi của bạn sau khi phân tích tài liệu')

    if user_question:
        answer = user_input(user_question)
        if answer:
            st.markdown("### 💬 Bot response:")
            st.write(answer)

    # Hiển thị lịch sử hỏi đáp
    if "chat_history" in st.session_state and st.session_state.chat_history:
        st.markdown("---")
        st.markdown("## 🕑 Chat history")
        for i, chat in enumerate(st.session_state.chat_history):
            st.markdown(f"**Q{i+1}:** {chat['question']}")
            st.markdown(f"**A{i+1}:** {chat['answer']}")
        #file word lưu tài liệu đã hỏi 
        if st.button("📄 Print answer"):
            word_file = export_chat_history_to_word(st.session_state.chat_history)
            st.download_button(
                label="📥 Tải file Word lịch sử chat",
                data=word_file,
                file_name="chat_history.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
            )
    if st.checkbox("📈 Hiện thống kê câu hỏi"):
        show_statistics()
        
    # Sidebar để tải file
    with st.sidebar:
        st.title('📁 Tải tài liệu')
        docs = st.file_uploader(
            'Tải tài liệu của bạn lên',
            accept_multiple_files=True,
            type=['pdf', 'docx', 'txt', 'pptx','png', 'jpg', 'jpeg', 'bmp']
        )

        if st.button('Phân tích tài liệu'):
            if not docs:
                st.error('Vui lòng tải tài liệu trước')
            else:
                with st.spinner('🔍 Đang xử lý tài liệu...'):
                    raw_text = get_text(docs)
                    if raw_text:
                        text_chunks = get_text_chunk(raw_text)
                        if text_chunks:
                            get_vector_store(text_chunks)
                        else:
                            st.error('Không thể chia nhỏ nội dung được trích xuất')
                    else:
                        st.error('Không thể đọc được nội dung tài liệu')
    
if __name__ == '__main__':
    main()